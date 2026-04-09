from __future__ import annotations

import json
import os
import re
import sys
from pathlib import Path

sys.path.insert(0, "/tmp/awareness_deps")

from pypdf import PdfReader  # type: ignore
from pptx import Presentation  # type: ignore


BASE_DIR = Path(
    "/Users/user/Air/SLG KB/SLG (internal)/ISO manual and policy/Information security policy"
)

FILES = [
    "20250331 Awareness training full company v1.1.pptx",
    "Internal Emergency Plan Keynius versie 1.2 16-4-2025.pdf",
    "Keynius Security Incident Response Plan v1.0.pdf",
    "Keynius SLG Information security policy v1.5.pdf",
    "Keynius Statement of Applicability v1.1 ENG.pdf",
    "Procedure safe use AI v1.0.pdf",
    "Procedure veilig gebruik AI v1.0.pdf",
    "SLG Beleidsregels Informatiebeveiliging Januari 2025.pdf",
    "Smart Locking Group security policies v1.2 januari 2025.pdf",
]


def normalize(text: str) -> str:
    text = text.replace("\xa0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return normalize("\n\n".join(pages))


def extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = [f"Slide {index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        slides.append("\n".join(parts))
    return normalize("\n\n".join(slides))


def main() -> None:
    output_dir = Path("build/extracted")
    output_dir.mkdir(parents=True, exist_ok=True)
    manifest = []

    for filename in FILES:
        source_path = BASE_DIR / filename
        if not source_path.exists():
            raise FileNotFoundError(source_path)

        if source_path.suffix.lower() == ".pdf":
            text = extract_pdf(source_path)
        elif source_path.suffix.lower() == ".pptx":
            text = extract_pptx(source_path)
        else:
            continue

        output_path = output_dir / f"{source_path.name}.txt"
        output_path.write_text(text, encoding="utf-8")
        manifest.append(
            {
                "source": str(source_path),
                "output": str(output_path.resolve()),
                "characters": len(text),
            }
        )

    manifest_path = output_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(manifest_path.resolve())


if __name__ == "__main__":
    main()
